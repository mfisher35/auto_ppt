import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


def add_pic_master_bottom_right(prs,pic_path):
  slide_master = prs.slide_master
  shapes = slide_master.shapes
  left = Inches(9.8)
  top = Inches(7.25)
  height = Inches(0.2)
  width = Inches(0.1)
  shape = shapes.add_picture(pic_path, left, top, width, height)

def change_all_title_fonts(prs,font):
  slide_master = prs.slide_master
  for i in range(1,len(slide_master.slide_layouts)):
    slide_master.slide_layouts[i].shapes[0].text_frame._set_font(font['family'],font['size'],font['bold'],font['italic'])
  for i in range(0,len(prs.slides)):
    try:
      prs.slides[i].shapes[0].text_frame._set_font(font['family'],font['size'],font['bold'],font['italic'])
    except:
      a = 1
