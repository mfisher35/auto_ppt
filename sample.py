import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
import helpers
import os

prs = Presentation('presentations/sample_inputs/Backend_Overview.pptx')
slide_master = prs.slide_master


helpers.add_logo_master_bottom_right(prs,'logo.png')
#shapes = slide_master.shapes
#print(dir(slide_master.shapes))
#slide_master.shapes.test()
#left = top = width = height = Inches(1.0)
#shape = shapes.add_picture(
#    'logo.png', left, top, width, height
#)

helpers.change_all_title_fonts(prs,{'family':'Calibri','size':36,'bold':0,'italic':0})
helpers.add_heading_line_breaks(prs)
helpers.align_all_titles_fonts(prs)
prs.save('new-file-name.pptx')
os.system('open new-file-name.pptx')
