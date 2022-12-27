import collections 
import collections.abc
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageEnhance


def add_heading_line_breaks(prs):
  slide_master = prs.slide_master
  for i in range(1,len(slide_master.slide_layouts)):
    shapes = slide_master.slide_layouts[i].shapes
    x1s = [Inches(0.4),Inches(0.4)]
    y1s = [Inches(1.00),Inches(6.95)]
    x2s = [Inches(9.6),Inches(9.6)]
    y2s = [Inches(1.00),Inches(6.95)]
    for i in range(0,len(x1s)):
      line = shapes.add_connector(MSO_CONNECTOR.STRAIGHT,x1s[i], y1s[i], x2s[i], y2s[i])
      shadow = line.shadow
      shadow.inherit = False
      line.line.fill.solid()
      line.line.fill.fore_color.rgb = RGBColor.from_string("FF0000")

def add_title_logo(prs,pic_path):
  im = Image.open(pic_path)
  pic_path = pic_path.replace(".",'-title.')
  enhancer = ImageEnhance.Sharpness(im)
  im2 = enhancer.enhance(0.01)
  enhancer = ImageEnhance.Contrast(im)
  im2 = enhancer.enhance(0.01)
  enhancer = ImageEnhance.Brightness(im)
  im2 = enhancer.enhance(3)
  im2.save(pic_path);
  shape = prs.slides[0].shapes.add_picture(pic_path, Inches(-2.72), Inches(-0.78), None, None)
  cursor_sp = prs.slides[0].shapes[0]._element
  cursor_sp.addprevious(shape._element)
  res_width, res_height =  _title_logo_proportional_width_height(shape.width,shape.height)
  shape.width = res_width
  shape.height = res_height
  shape.left = Inches(5) - shape.width 
  shape.top = Inches(0)
 

def add_logo_master_bottom_right(prs,pic_path):
  slide_master = prs.slide_master
  shapes = slide_master.shapes
  left = Inches(9.6)
  top = Inches(6.95)
  height = None
  width = None
  shape = shapes.add_picture(pic_path, left, top, width, height)
  res_width, res_height =  _get_proportional_width_height(shape.width,shape.height)
  shape.width = res_width
  shape.height = res_height
  shape.left = shape.left - shape.width 
  shape.top = shape.top + int(shape.height*0.75)
 
def change_all_title_fonts(prs,font):
  slide_master = prs.slide_master
  for i in range(1,len(slide_master.slide_layouts)):
    slide_master.slide_layouts[i].shapes[0].text_frame._set_font(font['family'],font['size'],font['bold'],font['italic'])
  for i in range(0,len(prs.slides)):
    try:
      prs.slides[i].shapes[0].text_frame._set_font(font['family'],font['size'],font['bold'],font['italic'])
    except:
      a = 1


def align_all_titles_fonts(prs):
  slide_master = prs.slide_master
  align_pos = {'left': Inches(0.4), 'width' : Inches(9.2), 'top' : Inches(0), 'height' : Inches(1)}

  for i in range(1,len(slide_master.slide_layouts)):
     shape = slide_master.slide_layouts[i].shapes[0]
     shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
     _apply_position(shape,align_pos)
  for i in range(1,len(prs.slides)):
    try:
      shape = prs.slides[i].shapes[0]
      shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
      _apply_position(shape,align_pos)
    except:
      a = 1


def _apply_position(obj,position):
   obj.left = position['left']
   obj.width = position['width']
   obj.top = position['top']
   obj.height = position['height']

 
def _get_proportional_width_height(width,height,target_width=0.25,target_height=0.25,method='min'):
 
  if (width > height and method=='min') or (width < height and method=='max'):
    ratio = height/width
    res_width = Inches(target_width)
    res_height = int(ratio*Inches(target_height))
  else:
    ratio = width/height
    res_height = Inches(target_height)
    res_width = int(ratio*Inches(target_width))
  return(res_width,res_height)

def _title_logo_proportional_width_height(width,height,slide_width=10,slide_height=7.5):
  slide_height = Inches(slide_height)
  slide_width = Inches(slide_width)
  ratio = slide_height/height
  res_height = int(ratio * height)
  res_width = int(ratio * width)
  return (res_width,res_height)
